
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

os.makedirs("artifacts", exist_ok=True)

# ── 색상 ──────────────────────────────────────────────────────
C_BG   = RGBColor(0x1A,0x1A,0x2E)
C_CARD = RGBColor(0x16,0x21,0x3E)
C_RED  = RGBColor(0xE9,0x4F,0x37)
C_YEL  = RGBColor(0xF5,0xA6,0x23)
C_WHT  = RGBColor(0xFF,0xFF,0xFF)
C_LGT  = RGBColor(0xCC,0xCC,0xCC)
C_GRY  = RGBColor(0x88,0x88,0x99)
C_GRN  = RGBColor(0x2E,0xCC,0x71)
C_BLU  = RGBColor(0x3A,0x9B,0xDC)
C_PUR  = RGBColor(0x9B,0x59,0xB6)
C_DRK  = RGBColor(0x1E,0x2A,0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]

def R(sl, l,t,w,h, f=None, lc=None, lw=1):
    s = sl.shapes.add_shape(1, Inches(l),Inches(t),Inches(w),Inches(h))
    if f: s.fill.solid(); s.fill.fore_color.rgb = f
    else: s.fill.background()
    if lc: s.line.color.rgb=lc; s.line.width=Pt(lw)
    else: s.line.fill.background()
    return s

def O(sl, l,t,w,h, f):
    s = sl.shapes.add_shape(9, Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=f; s.line.fill.background()
    return s

def T(sl, text, l,t,w,h, sz=13, bold=False, color=C_WHT, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=align
    r = p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.color.rgb=color; r.font.name="맑은 고딕"
    return tb

def BG(sl): R(sl,0,0,13.33,7.5,f=C_BG)
def BAR(sl,c=C_RED): R(sl,0,0,0.12,7.5,f=c)
def HDR(sl): R(sl,0,0,13.33,1.1,f=C_CARD)
def DIV(sl,l,t,w,c=C_RED): R(sl,l,t,w,0.04,f=c)

# ═══════════════════════════════════════════
# 슬라이드 1 – 표지
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl)
O(sl, 8.5,-1.2,5.5,5.5, RGBColor(0x22,0x22,0x44))
O(sl,-2.0, 4.0,4.5,4.5, RGBColor(0x1E,0x1E,0x3A))

T(sl,"🍽️  강남역 맛집 가이드",0.5,1.5,10,1.1,sz=40,bold=True)
T(sl,"2024-2025 현지인 추천 베스트 맛집 완벽 정리",0.5,2.8,10,0.6,sz=18,color=C_YEL)
DIV(sl,0.5,3.55,6.5)
T(sl,"구글·네이버 평점 및 현지인 리뷰를 기반으로 엄선한\n강남역 주변 인기 맛집을 카테고리별로 정리했습니다.",
   0.5,3.7,9,0.9,sz=13,color=C_LGT)

tags=[("🇻🇳 베트남",0),("🍝 양식",1),("🥩 스테이크",2),
      ("🍳 오므라이스",3),("🍲 한식",4),("🍜 일식",5)]
for tag,i in tags:
    c,row=i%3,i//3
    R(sl,0.5+c*2.2,5.0+row*0.6,2.0,0.45,f=C_CARD,lc=C_RED,lw=1)
    T(sl,tag,0.6+c*2.2,5.05+row*0.6,1.8,0.38,sz=11,align=PP_ALIGN.CENTER)

T(sl,"GANGNAM\nRESTAURANT\nGUIDE",9.8,2.2,3.2,2.5,sz=24,bold=True,
  color=RGBColor(0x2A,0x2A,0x4E),align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# 슬라이드 2 – 목차
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl); HDR(sl)
T(sl,"목  차",0.4,0.2,5,0.7,sz=30,bold=True)
DIV(sl,0.4,1.05,12.5)

items=[
    ("01","강남역 맛집 개요","강남역 상권 소개 및 맛집 선정 기준"),
    ("02","땀땀 – 베트남 쌀국수","소곱창 쌀국수의 원조, 줄 서는 맛집"),
    ("03","마녀주방 – 이색 레스토랑","365일 할로윈 컨셉 데이트 맛집"),
    ("04","미도인 – 스테이크 덮밥","한정 판매 400g 스테이크 덮밥"),
    ("05","을지다락 – 오므라이스","엘레강스한 분위기의 오므라이스 전문점"),
    ("06","강남진해장 – 곱창전골","5년 연속 블루리본 선정, 24시간 운영"),
    ("07","기타 추천 맛집","닭갈비·한정식·돈카츠·훠궈 등"),
    ("08","방문 꿀팁 & 총정리","예약 팁, 웨이팅 정보, 추천 코스"),
]
for i,(num,title,desc) in enumerate(items):
    col,row=i//4,i%4
    lx,ty=0.4+col*6.5,1.2+row*1.5
    R(sl,lx,ty,6.1,1.3,f=C_CARD,lc=C_YEL,lw=0.5)
    R(sl,lx,ty,0.65,1.3,f=C_RED)
    T(sl,num,lx+0.05,ty+0.35,0.55,0.55,sz=16,bold=True,align=PP_ALIGN.CENTER)
    T(sl,title,lx+0.75,ty+0.1,5.2,0.5,sz=14,bold=True,color=C_YEL)
    T(sl,desc, lx+0.75,ty+0.65,5.2,0.5,sz=11,color=C_LGT)

# ═══════════════════════════════════════════
# 슬라이드 3 – 강남역 맛집 개요
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl); HDR(sl)
T(sl,"01  강남역 맛집 개요",0.3,0.2,9,0.7,sz=24,bold=True)
T(sl,"OVERVIEW",10.5,0.28,2.5,0.5,sz=12,color=C_GRY,align=PP_ALIGN.RIGHT)

stats=[
    ("🏙️","서울 최대 상권","하루 유동인구\n약 20만 명"),
    ("🍴","다양한 음식 장르","한식·양식·일식·중식\n아시안 퓨전 등"),
    ("⭐","검증된 맛집","구글·네이버 평점\n4.0 이상 엄선"),
]
for i,(icon,title,desc) in enumerate(stats):
    lx=0.4+i*4.2
    R(sl,lx,1.25,3.9,2.0,f=C_CARD,lc=C_RED,lw=1)
    T(sl,icon, lx+0.15,1.35,0.8,0.8,sz=26,align=PP_ALIGN.CENTER)
    T(sl,title,lx+0.15,2.0, 3.6,0.45,sz=13,bold=True,color=C_YEL)
    T(sl,desc, lx+0.15,2.5, 3.6,0.65,sz=11,color=C_LGT)

T(sl,"📋  맛집 선정 기준",0.4,3.5,6,0.45,sz=15,bold=True,color=C_YEL)
criteria=[
    "✅  구글·네이버 양대 포털 평점 4.0 이상",
    "✅  블로그·방문자 리뷰 500건 이상",
    "✅  TV 방송 소개 또는 블루리본 선정 이력",
    "✅  현지인 재방문율 높은 가성비 맛집 우선",
]
for i,c in enumerate(criteria):
    T(sl,c,0.5,4.05+i*0.52,6,0.45,sz=12,color=C_LGT)

R(sl,7.0,3.5,5.9,3.6,f=C_CARD,lc=C_BLU,lw=1)
T(sl,"📍  강남역 출구별 맛집 밀집 구역",7.1,3.6,5.7,0.45,sz=13,bold=True,color=C_BLU)
exits=[
    ("6번 출구","안동국시, 을밀대 냉면"),
    ("9번 출구","포490 베트남 쌀국수"),
    ("10번 출구","정돈 돈카츠, 흑우점"),
    ("11번 출구","땀땀, 을지다락, 마녀주방"),
    ("12번 출구","도원참치, 미도인"),
]
for i,(ex,rest) in enumerate(exits):
    ty=4.15+i*0.52
    R(sl,7.1,ty,1.3,0.38,f=C_RED)
    T(sl,ex, 7.1,ty,1.3,0.38,sz=9,bold=True,align=PP_ALIGN.CENTER)
    T(sl,rest,8.5,ty+0.02,4.3,0.38,sz=11,color=C_LGT)

# ═══════════════════════════════════════════
# 슬라이드 4 – 땀땀
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl); HDR(sl)
T(sl,"02  땀땀  –  베트남 쌀국수 맛집",0.3,0.2,10,0.7,sz=24,bold=True)
T(sl,"강남역 11번 출구",10.0,0.28,3.0,0.5,sz=12,color=C_GRY,align=PP_ALIGN.RIGHT)

R(sl,0.3,1.2,7.5,5.9,f=C_CARD)
T(sl,"🍜  매장 소개",0.5,1.3,7,0.45,sz=15,bold=True,color=C_YEL)
intro=("강남역 11번 출구 근처에 위치한 베트남 쌀국수 전문점으로,\n"
       "베트남 전통 궁중 보양식 레시피를 바탕으로 소고기와 사골을\n"
       "24시간 우려낸 진한 국물이 일품입니다.\n\n"
       "생방송투데이·생방송오늘저녁·생생정보 등 다수 TV 방송에 소개된\n"
       "줄 서서 먹는 강남역 대표 쌀국수 맛집입니다.")
T(sl,intro,0.5,1.85,7.1,1.8,sz=12,color=C_LGT)
DIV(sl,0.5,3.7,7.0,c=C_YEL)
T(sl,"📋  대표 메뉴",0.5,3.85,7,0.4,sz=14,bold=True,color=C_YEL)

menus=[
    ("프리미엄 보양 쌀국수","17,000원","사골 24시간 우린 진한 국물"),
    ("매운 소곱창 쌀국수",  "16,000원","시그니처 메뉴, 매콤한 맛"),
    ("양지 쌀국수",         "11,000원","담백하고 깔끔한 국물"),
    ("우삼겹 쌀국수",       "12,000원","부드러운 우삼겹 토핑"),
]
for i,(name,price,desc) in enumerate(menus):
    ty=4.35+i*0.55
    R(sl,0.5,ty,7.1,0.48,f=C_DRK)
    T(sl,f"• {name}",0.65,ty+0.06,3.5,0.38,sz=12)
    R(sl,4.2,ty+0.06,1.3,0.35,f=C_RED)
    T(sl,price,4.2,ty+0.06,1.3,0.35,sz=11,bold=True,align=PP_ALIGN.CENTER)
    T(sl,desc,5.6,ty+0.08,2.0,0.35,sz=10,color=C_GRY)

R(sl,8.0,1.2,5.0,5.9,f=C_CARD,lc=C_RED,lw=1.5)
T(sl,"📌  매장 정보",8.2,1.3,4.6,0.45,sz=14,bold=True,color=C_RED)
info=[
    ("📍 위치",    "서울 강남구 강남대로98길 12-5"),
    ("⏰ 영업시간","11:00~22:00 (라스트오더 20:30)"),
    ("📞 연락처",  "02-554-8892"),
    ("⭐ 구글 평점","4.1 / 5.0"),
    ("💡 특징",    "줄 서는 맛집 / TV 방송 다수 소개"),
]
for i,(label,val) in enumerate(info):
    ty=1.9+i*0.85
    R(sl,8.1,ty,4.7,0.72,f=C_DRK)
    T(sl,label,8.2,ty+0.04,1.8,0.3,sz=10,bold=True,color=C_YEL)
    T(sl,val,  8.2,ty+0.34,4.5,0.32,sz=11)

# ═══════════════════════════════════════════
# 슬라이드 5 – 마녀주방
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl,c=C_PUR); HDR(sl)
T(sl,"03  마녀주방 강남점  –  이색 테마 레스토랑",0.3,0.2,11,0.7,sz=24,bold=True)
T(sl,"강남역 11번 출구 B1",10.5,0.28,2.6,0.5,sz=12,color=C_GRY,align=PP_ALIGN.RIGHT)

R(sl,0.3,1.2,7.5,5.9,f=C_CARD)
T(sl,"🧙‍♀️  매장 소개",0.5,1.3,7,0.45,sz=15,bold=True,color=C_PUR)
intro=("365일 할로윈을 즐길 수 있는 이색 테마 레스토랑으로,\n"
       "독특한 인테리어와 함께 맛있는 양식을 즐길 수 있어\n"
       "데이트 코스 및 특별한 날 방문하기 좋은 곳입니다.\n\n"
       "구글 평점 4.2를 기록하며 강남역 맛집 TOP 2에 선정된\n"
       "인스타그램 감성 가득한 레스토랑입니다.")
T(sl,intro,0.5,1.85,7.1,1.8,sz=12,color=C_LGT)
DIV(sl,0.5,3.7,7.0,c=C_PUR)
T(sl,"📋  인기 메뉴",0.5,3.85,7,0.4,sz=14,bold=True,color=C_PUR)

menus=[
    ("넓적다리 스테이크","시그니처","육즙 풍부한 스테이크"),
    ("링거 칵테일",      "이색 음료","할로윈 컨셉 칵테일"),
    ("포테이토 피자",    "인기 사이드","바삭한 감자 토핑 피자"),
    ("까르보나라 파스타","양식 메뉴","크리미한 소스의 파스타"),
]
for i,(name,tag,desc) in enumerate(menus):
    ty=4.35+i*0.55
    R(sl,0.5,ty,7.1,0.48,f=C_DRK)
    T(sl,f"• {name}",0.65,ty+0.06,3.0,0.38,sz=12)
    R(sl,3.7,ty+0.06,1.5,0.35,f=C_PUR)
    T(sl,tag,3.7,ty+0.06,1.5,0.35,sz=10,bold=True,align=PP_ALIGN.CENTER)
    T(sl,desc,5.3,ty+0.08,2.3,0.35,sz=10,color=C_GRY)

R(sl,8.0,1.2,5.0,5.9,f=C_CARD,lc=C_PUR,lw=1.5)
T(sl,"📌  매장 정보",8.2,1.3,4.6,0.45,sz=14,bold=True,color=C_PUR)
info=[
    ("📍 위치",    "서울 강남구 강남대로94길 9, B1층"),
    ("⏰ 영업시간","12:00~22:30 (브레이크 16:00~17:30)"),
    ("📞 연락처",  "070-4240-1116"),
    ("⭐ 구글 평점","4.2 / 5.0"),
    ("💡 특징",    "365일 할로윈 / 데이트 코스 인기"),
]
for i,(label,val) in enumerate(info):
    ty=1.9+i*0.85
    R(sl,8.1,ty,4.7,0.72,f=C_DRK)
    T(sl,label,8.2,ty+0.04,1.8,0.3,sz=10,bold=True,color=C_PUR)
    T(sl,val,  8.2,ty+0.34,4.5,0.32,sz=11)

# ═══════════════════════════════════════════
# 슬라이드 6 – 미도인 & 을지다락
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl,c=C_GRN); HDR(sl)
T(sl,"04 & 05  미도인  /  을지다락",0.3,0.2,10,0.7,sz=24,bold=True)
T(sl,"스테이크 덮밥 & 오므라이스",9.5,0.28,3.6,0.5,sz=12,color=C_GRY,align=PP_ALIGN.RIGHT)

# 미도인 (왼쪽)
R(sl,0.3,1.2,6.1,5.9,f=C_CARD)
R(sl,0.3,1.2,6.1,0.5,f=C_GRN)
T(sl,"🥩  미도인 강남  –  스테이크 덮밥의 명가",0.45,1.27,5.8,0.38,sz=12,bold=True)
T(sl,"강남역 CGV 근처에 위치한 미도인은 하루 7인분 한정 판매하는\n400g 스테이크 덮밥으로 유명한 맛집입니다.\n스테이크와 한식의 절묘한 조화가 특징입니다.",
   0.45,1.85,5.8,1.1,sz=11,color=C_LGT)
T(sl,"대표 메뉴",0.45,3.05,5.8,0.35,sz=12,bold=True,color=C_GRN)
menus_m=[("스테이크 덮밥","16,800원"),("미도인 스테이크 덮밥","10,800원"),
         ("9첩 반상","15,300원"),("가정식 등심스테이크","17,600원")]
for i,(name,price) in enumerate(menus_m):
    ty=3.5+i*0.48
    R(sl,0.45,ty,5.8,0.42,f=C_DRK)
    T(sl,f"• {name}",0.6,ty+0.06,3.5,0.32,sz=11)
    R(sl,4.1,ty+0.05,1.9,0.32,f=C_GRN)
    T(sl,price,4.1,ty+0.05,1.9,0.32,sz=10,bold=True,align=PP_ALIGN.CENTER)
for i,(icon,val) in enumerate([("📍","강남구 강남대로102길 16"),("⏰","11:20~21:00 (브레이크 15~17시)"),("⭐","구글 평점 4.0+")]):
    T(sl,f"{icon} {val}",0.45,5.45+i*0.38,5.8,0.35,sz=10,color=C_GRY)

# 을지다락 (오른쪽)
R(sl,6.7,1.2,6.3,5.9,f=C_CARD)
R(sl,6.7,1.2,6.3,0.5,f=C_BLU)
T(sl,"🍳  을지다락 강남  –  오므라이스 맛집",6.85,1.27,6.0,0.38,sz=12,bold=True)
T(sl,"강남역 11번 출구 근처에 위치한 오므라이스 전문 레스토랑으로,\n엘레강스한 인테리어와 촉촉한 에그 스크램블이 특징입니다.\n소개팅·데이트 코스로 인기 높은 강남역 레스토랑입니다.",
   6.85,1.85,6.0,1.1,sz=11,color=C_LGT)
T(sl,"대표 메뉴",6.85,3.05,6.0,0.35,sz=12,bold=True,color=C_BLU)
menus_e=[("다락 오므라이스","16,000원"),("다락 로제파스타","18,000원"),
         ("참목살 스테이크","22,000원"),("가츠산도","12,000원")]
for i,(name,price) in enumerate(menus_e):
    ty=3.5+i*0.48
    R(sl,6.85,ty,6.0,0.42,f=C_DRK)
    T(sl,f"• {name}",7.0,ty+0.06,3.8,0.32,sz=11)
    R(sl,10.7,ty+0.05,1.9,0.32,f=C_BLU)
    T(sl,price,10.7,ty+0.05,1.9,0.32,sz=10,bold=True,align=PP_ALIGN.CENTER)
for i,(icon,val) in enumerate([("📍","강남구 강남대로 96길 22 2층"),("⏰","11:30~21:10 (브레이크 15:10~16:30)"),("⭐","구글 평점 4.3")]):
    T(sl,f"{icon} {val}",6.85,5.45+i*0.38,6.0,0.35,sz=10,color=C_GRY)

# ═══════════════════════════════════════════
# 슬라이드 7 – 강남진해장
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl,c=C_YEL); HDR(sl)
T(sl,"06  강남진해장  –  곱창전골 맛집",0.3,0.2,10,0.7,sz=24,bold=True)
T(sl,"24시간 운영 | 5년 연속 블루리본",9.0,0.28,4.1,0.5,sz=12,color=C_YEL,align=PP_ALIGN.RIGHT)

R(sl,0.3,1.2,7.5,5.9,f=C_CARD)
T(sl,"🍲  매장 소개",0.5,1.3,7,0.45,sz=15,bold=True,color=C_YEL)
intro=("5년 연속 블루리본 서베이에 선정된 강남진해장은\n"
       "곱창전골로 유명한 24시간 운영 맛집입니다.\n\n"
       "진한 육수와 신선한 곱창이 어우러진 전골 요리로\n"
       "강남역 직장인들의 회식 장소로 큰 인기를 끌고 있습니다.\n\n"
       "야식·해장·회식 등 어떤 목적으로도 방문하기 좋은\n"
       "강남역 대표 한식 맛집입니다.")
T(sl,intro,0.5,1.85,7.1,2.2,sz=12,color=C_LGT)
DIV(sl,0.5,4.1,7.0,c=C_YEL)
T(sl,"📋  인기 메뉴",0.5,4.25,7,0.4,sz=14,bold=True,color=C_YEL)
menus=[
    ("곱창전골",    "65,000원","시그니처 메뉴 (2인 기준)"),
    ("양선지해장국","12,000원","진한 국물의 해장 메뉴"),
    ("사골곰탕",    "12,000원","담백하고 깊은 맛"),
    ("모듬전골",    "70,000원","다양한 재료의 모듬 전골"),
]
for i,(name,price,desc) in enumerate(menus):
    ty=4.75+i*0.52
    R(sl,0.5,ty,7.1,0.46,f=C_DRK)
    T(sl,f"• {name}",0.65,ty+0.06,3.0,0.36,sz=12)
    R(sl,3.7,ty+0.06,1.5,0.33,f=C_YEL)
    T(sl,price,3.7,ty+0.06,1.5,0.33,sz=10,bold=True,color=C_BG,align=PP_ALIGN.CENTER)
    T(sl,desc,5.3,ty+0.08,2.3,0.33,sz=10,color=C_GRY)

R(sl,8.0,1.2,5.0,5.9,f=C_CARD,lc=C_YEL,lw=1.5)
T(sl,"📌  매장 정보",8.2,1.3,4.6,0.45,sz=14,bold=True,color=C_YEL)
info=[
    ("📍 위치",    "서울 강남구 테헤란로5길 11"),
    ("⏰ 영업시간","24시간 연중무휴"),
    ("📞 연락처",  "02-557-2662"),
    ("⭐ 구글 평점","4.2 / 5.0"),
]
for i,(label,val) in enumerate(info):
    ty=1.9+i*0.75
    R(sl,8.1,ty,4.7,0.65,f=C_DRK)
    T(sl,label,8.2,ty+0.04,1.8,0.28,sz=10,bold=True,color=C_YEL)
    T(sl,val,  8.2,ty+0.32,4.5,0.28,sz=11)
R(sl,8.1,5.0,4.7,1.8,f=RGBColor(0x2A,0x1A,0x00),lc=C_YEL,lw=1)
T(sl,"🏆  수상 & 인증 내역",8.2,5.08,4.5,0.38,sz=12,bold=True,color=C_YEL)
for i,a in enumerate(["🎖️ 블루리본 서베이 5년 연속 선정","📺 TV 방송 다수 소개","💯 네이버 플레이스 상위 노출"]):
    T(sl,a,8.2,5.55+i*0.38,4.5,0.35,sz=10,color=C_LGT)

# ═══════════════════════════════════════════
# 슬라이드 8 – 기타 추천 맛집
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl,c=C_BLU); HDR(sl)
T(sl,"07  기타 추천 맛집",0.3,0.2,10,0.7,sz=24,bold=True)
T(sl,"다양한 장르의 강남역 맛집",9.5,0.28,3.6,0.5,sz=12,color=C_GRY,align=PP_ALIGN.RIGHT)

others=[
    ("🍗","장인닭갈비",   "달달하고 맵지 않은 닭갈비\n12,000원~14,000원","강남구 테헤란로1길 19","11:00~24:00",C_RED),
    ("🍱","노랑저고리 한정식","상다리 부러지는 한정식\n점심 19,000원~","서초구 서초대로73길 9 5층","11:00~21:30",C_GRN),
    ("🥩","정돈 강남점",  "바삭한 일본식 돈카츠\n등심돈카츠 16,000원","신논현역 5번 출구 근처","오픈 직후 방문 추천",C_YEL),
    ("🍜","하이디라오",   "중국 훠궈의 정수\n탕 선택 폭 넓음","신논현역 인근","예약 필수",RGBColor(0xFF,0x6B,0x6B)),
    ("🍝","고에몬 강남점","일본식 파스타 전문점\n까르보나라 13,000원~","강남역 인근","점심·저녁 운영",C_BLU),
    ("🥘","뱅뱅막국수",  "들기름막국수 숨은 맛집\n구글 평점 4.5","뱅뱅사거리 인근","점심 영업",C_PUR),
]
for i,(icon,name,desc,addr,hours,color) in enumerate(others):
    col,row=i%3,i//3
    lx,ty=0.3+col*4.3,1.25+row*2.9
    R(sl,lx,ty,4.0,2.6,f=C_CARD,lc=color,lw=1.5)
    R(sl,lx,ty,4.0,0.5,f=color)
    T(sl,f"{icon}  {name}",lx+0.1,ty+0.08,3.8,0.38,sz=13,bold=True)
    T(sl,desc, lx+0.1,ty+0.6, 3.8,0.7, sz=11,color=C_LGT)
    DIV(sl,lx+0.1,ty+1.35,3.7,c=color)
    T(sl,f"📍 {addr}", lx+0.1,ty+1.45,3.8,0.35,sz=9,color=C_GRY)
    T(sl,f"⏰ {hours}",lx+0.1,ty+1.85,3.8,0.35,sz=9,color=C_GRY)

# ═══════════════════════════════════════════
# 슬라이드 9 – 방문 꿀팁 & 총정리
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl); BAR(sl); HDR(sl)
T(sl,"08  방문 꿀팁 & 총정리",0.3,0.2,10,0.7,sz=24,bold=True)

R(sl,0.3,1.2,6.1,5.9,f=C_CARD)
T(sl,"💡  방문 꿀팁",0.5,1.3,5.8,0.45,sz=15,bold=True,color=C_YEL)
tips=[
    ("⏰ 방문 시간",  "평일 점심(12~13시)은 웨이팅이 길어요.\n오픈 직후(11~12시) 또는 14시 이후 방문 추천!"),
    ("📱 예약 방법",  "네이버 예약 가능한 곳은 미리 예약하세요.\n특히 주말·공휴일은 필수입니다."),
    ("🚇 교통 정보",  "강남역 2호선 이용 시 출구별 맛집 위치 확인!\n11번 출구 주변에 맛집이 가장 많습니다."),
    ("💰 가성비 팁",  "점심 특선 메뉴를 활용하면 저렴하게 즐길 수 있어요.\n대부분 10,000~18,000원 선입니다."),
    ("👥 인원 구성",  "2~4인 방문 시 다양한 메뉴를 나눠 먹기 좋아요.\n1인 방문도 가능한 곳이 많습니다."),
]
for i,(title,content) in enumerate(tips):
    ty=1.9+i*1.0
    R(sl,0.4,ty,5.9,0.88,f=C_DRK)
    R(sl,0.4,ty,0.06,0.88,f=C_RED)
    T(sl,title,  0.55,ty+0.05,5.6,0.3, sz=11,bold=True,color=C_YEL)
    T(sl,content,0.55,ty+0.38,5.6,0.45,sz=10,color=C_LGT)

R(sl,6.7,1.2,6.3,5.9,f=C_CARD)
T(sl,"📊  맛집 총정리 요약",6.9,1.3,6.0,0.45,sz=15,bold=True,color=C_YEL)
headers=["맛집명","장르","가격대","평점"]
col_w=[1.8,1.4,1.4,0.9]
col_s=[6.85,8.65,10.05,11.45]
for j,(h,w,s) in enumerate(zip(headers,col_w,col_s)):
    R(sl,s,1.9,w-0.05,0.38,f=C_RED)
    T(sl,h,s+0.05,1.92,w-0.1,0.34,sz=10,bold=True,align=PP_ALIGN.CENTER)
rows=[
    ("땀땀",      "베트남 쌀국수","11,000~17,000","⭐4.1"),
    ("마녀주방",  "양식/테마",    "15,000~25,000","⭐4.2"),
    ("미도인",    "스테이크 덮밥","10,800~17,600","⭐4.0+"),
    ("을지다락",  "오므라이스",   "12,000~22,000","⭐4.3"),
    ("강남진해장","곱창전골",     "12,000~70,000","⭐4.2"),
    ("장인닭갈비","닭갈비",       "12,000~14,000","⭐4.0+"),
    ("정돈 강남", "돈카츠",       "13,500~45,800","⭐4.5"),
    ("하이디라오","훠궈",         "1인 30,000+",  "⭐4.4"),
]
row_colors=[C_DRK, RGBColor(0x22,0x2E,0x50)]
for i,row_data in enumerate(rows):
    ty=2.38+i*0.52
    R(sl,6.85,ty,5.5,0.46,f=row_colors[i%2])
    for j,(val,w,s) in enumerate(zip(row_data,col_w,col_s)):
        fc=C_YEL if j==3 else C_WHT
        T(sl,val,s+0.05,ty+0.08,w-0.1,0.32,sz=9,color=fc,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# 슬라이드 10 – 마무리
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(BL)
BG(sl)
O(sl, 8.0,-1.5,6.0,6.0, RGBColor(0x22,0x22,0x44))
O(sl,-2.0, 3.5,5.0,5.0, RGBColor(0x1E,0x1E,0x3A))
R(sl,0,0,0.12,7.5,f=C_RED)

T(sl,"🍽️",5.5,1.2,2.0,1.5,sz=60,align=PP_ALIGN.CENTER)
T(sl,"강남역 맛집 가이드",2.0,2.8,9.0,1.0,sz=36,bold=True,align=PP_ALIGN.CENTER)
DIV(sl,3.5,3.95,6.0)
T(sl,"맛있는 식사와 함께 즐거운 강남 나들이 되세요! 😊",
   1.5,4.15,10.0,0.7,sz=17,color=C_LGT,align=PP_ALIGN.CENTER)

tags_f=["#강남역맛집","#서울맛집","#데이트코스","#회식장소","#점심맛집","#2025맛집"]
for i,tag in enumerate(tags_f):
    lx=1.0+i*1.9
    R(sl,lx,5.2,1.75,0.42,f=C_CARD,lc=C_RED,lw=0.8)
    T(sl,tag,lx+0.05,5.24,1.65,0.35,sz=10,color=C_YEL,align=PP_ALIGN.CENTER)

T(sl,"Powered by Tavily Search  |  2024-2025 강남역 맛집 정보",
   1.5,6.8,10.0,0.45,sz=10,color=C_GRY,align=PP_ALIGN.CENTER)

# ─── 저장 ────────────────────────────────────────────────────
out = "artifacts/강남역_맛집_가이드.pptx"
prs.save(out)
print(f"✅ 저장 완료: {out}")
print(f"📊 총 슬라이드 수: {len(prs.slides)}장")
