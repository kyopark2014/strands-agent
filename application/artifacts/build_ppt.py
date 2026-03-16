
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.makedirs("artifacts", exist_ok=True)

# 색상
C_BG     = RGBColor(0x0D,0x1B,0x2A)
C_ACCENT = RGBColor(0xFF,0x99,0x00)
C_ACCENT2= RGBColor(0x00,0xA8,0xE8)
C_WHITE  = RGBColor(0xFF,0xFF,0xFF)
C_LIGHT  = RGBColor(0xD0,0xD8,0xE8)
C_CARD   = RGBColor(0x16,0x2A,0x42)
C_CARD2  = RGBColor(0x1E,0x35,0x52)
C_GREEN  = RGBColor(0x2E,0xCC,0x71)
C_PURPLE = RGBColor(0x9B,0x59,0xB6)
C_RED    = RGBColor(0xE7,0x4C,0x3C)
C_TEAL   = RGBColor(0x1A,0xBC,0x9C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def ns():
    return prs.slides.add_slide(prs.slide_layouts[6])
def bg(sl,c=C_BG):
    b=sl.background; b.fill.solid(); b.fill.fore_color.rgb=c
def rc(sl,l,t,w,h,fc,lc=None,lw=Pt(1)):
    sh=sl.shapes.add_shape(1,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if lc: sh.line.color.rgb=lc; sh.line.width=lw
    else:  sh.line.fill.background()
    return sh
def tx(sl,text,l,t,w,h,sz=Pt(14),bold=False,color=C_WHITE,
       align=PP_ALIGN.LEFT,italic=False):
    tb=sl.shapes.add_textbox(l,t,w,h)
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=sz; r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=color
    return tb
def ml(sl,lines,l,t,w,h,dsz=Pt(13),dc=C_WHITE,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(l,t,w,h)
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(lines):
        if isinstance(item,str): item=(item,False,dc,dsz)
        tx2,bld,clr,sz=item
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        r=p.add_run(); r.text=tx2
        r.font.size=sz; r.font.bold=bld; r.font.color.rgb=clr
    return tb
def topbar(sl):
    rc(sl,0,0,W,Inches(0.07),C_ACCENT)
def snum(sl,n,total=9):
    tx(sl,f"{n} / {total}",W-Inches(1.3),H-Inches(0.38),
       Inches(1.1),Inches(0.3),sz=Pt(10),color=C_LIGHT,align=PP_ALIGN.RIGHT)
def stitle(sl,title,sub=None):
    tx(sl,title,Inches(0.45),Inches(0.55),Inches(12.4),Inches(0.65),
       sz=Pt(28),bold=True,color=C_WHITE)
    if sub:
        tx(sl,sub,Inches(0.45),Inches(1.18),Inches(11),Inches(0.4),
           sz=Pt(13),color=C_LIGHT)
    rc(sl,Inches(0.45),Inches(1.55),Inches(12.4),Inches(0.04),C_ACCENT2)
def slabel(sl,label):
    tx(sl,label,Inches(0.45),Inches(0.18),Inches(5),Inches(0.35),
       sz=Pt(11),color=C_ACCENT,bold=True)

# ══════════════════════════════════════════════════════════════════
# S1 표지
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s)
rc(s,W*0.58,0,W*0.42,H,RGBColor(0x10,0x22,0x38))
rc(s,W*0.76,0,W*0.24,H,RGBColor(0x13,0x28,0x44))
topbar(s)
rc(s,Inches(1.0),Inches(1.7),Inches(0.07),Inches(3.6),C_ACCENT)
tx(s,"AWS  |  Amazon Bedrock",Inches(1.2),Inches(0.2),Inches(6),Inches(0.5),
   sz=Pt(13),color=C_ACCENT,bold=True)
tx(s,"AgentCore",Inches(1.2),Inches(1.8),Inches(8),Inches(1.5),
   sz=Pt(64),bold=True,color=C_WHITE)
tx(s,"AI 에이전트 구축 · 배포 · 운영을 위한\n엔터프라이즈급 완전 관리형 플랫폼",
   Inches(1.2),Inches(3.35),Inches(7.5),Inches(1.1),sz=Pt(20),color=C_LIGHT)
rc(s,Inches(1.2),Inches(4.62),Inches(5.6),Inches(0.52),C_ACCENT)
tx(s,"Any Framework  ·  Any Model  ·  Zero Infra Hassle",
   Inches(1.2),Inches(4.62),Inches(5.6),Inches(0.52),
   sz=Pt(13),bold=True,color=C_BG,align=PP_ALIGN.CENTER)
tx(s,"GA: Oct 13, 2025   |   Latest Update: Dec 2, 2025",
   Inches(1.2),Inches(5.38),Inches(6),Inches(0.4),sz=Pt(11),color=C_LIGHT)
tx(s,"🤖",Inches(8.2),Inches(2.0),Inches(3.8),Inches(3.0),
   sz=Pt(120),align=PP_ALIGN.CENTER)
snum(s,1)

# ══════════════════════════════════════════════════════════════════
# S2 목차
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"TABLE OF CONTENTS")
stitle(s,"목차","Amazon Bedrock AgentCore 전체 구성")
toc=[
    ("01","AgentCore 개요","배경 · 정의 · 출시 일정"),
    ("02","핵심 서비스 구성","6가지 모듈 아키텍처"),
    ("03","AgentCore Runtime","실행 환경 · 세션 · 스트리밍"),
    ("04","Gateway & Identity","도구 연결 · 인증·인가"),
    ("05","Memory & Observability","메모리 · 모니터링"),
    ("06","최신 기능 (2025.12)","Policy · Evaluations"),
    ("07","지원 프레임워크 & 모델","오픈 아키텍처 전략"),
    ("08","엔터프라이즈 보안","VPC · PrivateLink · IAM"),
    ("09","활용 사례 & 요약","Use Cases · 핵심 가치"),
]
cols=[(Inches(0.5),Inches(1.75)),(Inches(4.55),Inches(1.75)),(Inches(8.6),Inches(1.75))]
for idx,(num,title,sub) in enumerate(toc):
    col=idx//3; row=idx%3
    lx,ty=cols[col]; ty=ty+row*Inches(1.55)
    cc=C_CARD if idx%2==0 else C_CARD2
    rc(s,lx,ty,Inches(3.7),Inches(1.35),cc,lc=C_ACCENT2,lw=Pt(0.8))
    rc(s,lx,ty,Inches(0.55),Inches(1.35),C_ACCENT)
    tx(s,num,lx,ty+Inches(0.35),Inches(0.55),Inches(0.65),
       sz=Pt(16),bold=True,color=C_BG,align=PP_ALIGN.CENTER)
    tx(s,title,lx+Inches(0.65),ty+Inches(0.15),Inches(2.95),Inches(0.5),
       sz=Pt(14),bold=True,color=C_WHITE)
    tx(s,sub,lx+Inches(0.65),ty+Inches(0.65),Inches(2.95),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
snum(s,2)

# ══════════════════════════════════════════════════════════════════
# S3 개요
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"01  OVERVIEW")
stitle(s,"AgentCore 개요","Amazon Bedrock AgentCore란 무엇인가?")
rc(s,Inches(0.45),Inches(1.75),Inches(5.8),Inches(5.3),C_CARD,lc=C_ACCENT2,lw=Pt(0.8))
rc(s,Inches(0.45),Inches(1.75),Inches(5.8),Inches(0.45),C_ACCENT2)
tx(s,"  정의 (What is AgentCore?)",Inches(0.45),Inches(1.75),Inches(5.8),Inches(0.45),
   sz=Pt(13),bold=True,color=C_BG)
ml(s,[
    ("Amazon Bedrock AgentCore는 어떤 프레임워크와 파운데이션 모델을 사용하더라도",False,C_LIGHT,Pt(13)),
    ("AI 에이전트를 안전하게 구축·배포·운영할 수 있는 완전 관리형 에이전틱 플랫폼입니다.",False,C_LIGHT,Pt(13)),
    ("",False,C_LIGHT,Pt(6)),
    ("인프라 관리 없이 에이전트가 도구와 데이터에 걸쳐 작업을 수행하고,",False,C_LIGHT,Pt(13)),
    ("프로덕션 환경에서 성능을 모니터링할 수 있습니다.",False,C_LIGHT,Pt(13)),
],Inches(0.6),Inches(2.3),Inches(5.5),Inches(1.5))
milestones=[
    (C_PURPLE,"2025 Preview","서비스 프리뷰 공개\nRuntime · Gateway · Identity"),
    (C_ACCENT,"Oct 13, 2025","정식 출시 (GA)\nVPC · PrivateLink · A2A 지원"),
    (C_GREEN, "Dec 2, 2025","Policy · Evaluations 추가\n에피소딕 메모리 · 양방향 스트리밍"),
]
for i,(clr,date,desc) in enumerate(milestones):
    ty=Inches(3.2)+i*Inches(0.95)
    rc(s,Inches(0.6),ty,Inches(0.35),Inches(0.35),clr)
    tx(s,date,Inches(1.1),ty,Inches(1.5),Inches(0.4),sz=Pt(12),bold=True,color=clr)
    tx(s,desc,Inches(2.7),ty,Inches(3.3),Inches(0.7),sz=Pt(11),color=C_LIGHT)
kv=[
    (C_ACCENT,"🚀","빠른 배포","에이전트를 분 단위로 생성·배포"),
    (C_ACCENT2,"🔒","엔터프라이즈 보안","VPC·IAM·OAuth 2.0 내장"),
    (C_GREEN,"⚡","고성능 실행","최대 8시간 장시간 실행 지원"),
    (C_PURPLE,"🔍","완전한 가시성","CloudWatch 기반 통합 모니터링"),
]
for i,(clr,icon,t2,d2) in enumerate(kv):
    col=i%2; row=i//2
    lx=Inches(6.7)+col*Inches(3.2)
    ty=Inches(1.75)+row*Inches(2.55)
    rc(s,lx,ty,Inches(3.0),Inches(2.3),C_CARD,lc=clr,lw=Pt(1.2))
    tx(s,icon,lx+Inches(0.15),ty+Inches(0.15),Inches(0.6),Inches(0.6),sz=Pt(24))
    tx(s,t2,lx+Inches(0.15),ty+Inches(0.75),Inches(2.7),Inches(0.45),
       sz=Pt(15),bold=True,color=clr)
    tx(s,d2,lx+Inches(0.15),ty+Inches(1.2),Inches(2.7),Inches(0.9),
       sz=Pt(12),color=C_LIGHT)
snum(s,3)

# ══════════════════════════════════════════════════════════════════
# S4 핵심 서비스 구성
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"02  ARCHITECTURE")
stitle(s,"핵심 서비스 구성","6가지 모듈형 서비스 – 독립적 또는 통합 사용 가능")
modules=[
    (C_ACCENT,  "🚀","Runtime",   "에이전트 실행 환경\n최대 8시간 · 자동 스케일링\nMCP 서버 · A2A 프로토콜"),
    (C_ACCENT2, "🔗","Gateway",   "API·Lambda를 MCP로 변환\n시맨틱 도구 검색\nSlack·Salesforce 연동"),
    (C_GREEN,   "🔐","Identity",  "OAuth 2.0 / OIDC 인증\n2LO · 3LO 지원\n최소 권한 접근 제어"),
    (C_PURPLE,  "🧠","Memory",    "단기·장기 메모리 관리\n에피소딕 메모리\n자체 관리 전략 지원"),
    (C_TEAL,    "📊","Observability","OTel 호환 모니터링\nCloudWatch 대시보드\nDatadog·LangSmith 연동"),
    (C_RED,     "🛠️","Tools",     "코드 인터프리터 내장\n브라우저 도구 제공\n확장 가능한 도구 생태계"),
]
for i,(clr,icon,name,desc) in enumerate(modules):
    col=i%3; row=i//3
    lx=Inches(0.45)+col*Inches(4.25)
    ty=Inches(1.75)+row*Inches(2.55)
    rc(s,lx,ty,Inches(4.0),Inches(2.35),C_CARD,lc=clr,lw=Pt(1.5))
    rc(s,lx,ty,Inches(4.0),Inches(0.45),clr)
    tx(s,icon,lx+Inches(0.12),ty+Inches(0.02),Inches(0.45),Inches(0.42),
       sz=Pt(18),align=PP_ALIGN.CENTER)
    tx(s,f"AgentCore {name}",lx+Inches(0.6),ty+Inches(0.05),Inches(3.3),Inches(0.38),
       sz=Pt(14),bold=True,color=C_BG)
    tx(s,desc,lx+Inches(0.18),ty+Inches(0.6),Inches(3.65),Inches(1.6),
       sz=Pt(12),color=C_LIGHT)
snum(s,4)

# ══════════════════════════════════════════════════════════════════
# S5 Runtime
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"03  RUNTIME")
stitle(s,"AgentCore Runtime","에이전트 실행 환경 – 인프라 없이 안전하게 실행")
features=[
    (C_ACCENT,"⏱️ 최대 8시간 실행","장시간 비동기 처리 지원\n복잡한 멀티스텝 에이전트 워크플로 처리 가능"),
    (C_ACCENT2,"🔄 자동 스케일링","동시 세션 수에 따라 자동 확장\n용량 계획 및 인프라 유지보수 불필요"),
    (C_GREEN,"🔒 완전한 세션 격리","각 사용자 세션 독립 실행\n보안 샌드박스 환경 보장"),
    (C_PURPLE,"📡 양방향 스트리밍","동시 듣기·응답 처리\n음성 에이전트 구현 지원"),
]
for i,(clr,title2,desc2) in enumerate(features):
    ty=Inches(1.75)+i*Inches(1.3)
    rc(s,Inches(0.45),ty,Inches(5.8),Inches(1.15),C_CARD,lc=clr,lw=Pt(1.2))
    rc(s,Inches(0.45),ty,Inches(0.08),Inches(1.15),clr)
    tx(s,title2,Inches(0.65),ty+Inches(0.1),Inches(5.4),Inches(0.4),
       sz=Pt(13),bold=True,color=clr)
    tx(s,desc2,Inches(0.65),ty+Inches(0.5),Inches(5.4),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
rc(s,Inches(6.7),Inches(1.75),Inches(6.2),Inches(5.3),C_CARD,lc=C_ACCENT2,lw=Pt(0.8))
rc(s,Inches(6.7),Inches(1.75),Inches(6.2),Inches(0.45),C_ACCENT2)
tx(s,"  배포 흐름 (Deployment Flow)",Inches(6.7),Inches(1.75),Inches(6.2),Inches(0.45),
   sz=Pt(13),bold=True,color=C_BG)
steps=[
    (C_ACCENT,"1","에이전트 코드 작성","LangGraph·CrewAI 등 프레임워크 사용\nHTTP 엔드포인트 추가"),
    (C_ACCENT2,"2","컨테이너 이미지 빌드","Amazon ECR에 이미지 푸시\nAgentCore SDK 활용 가능"),
    (C_GREEN,"3","Runtime 생성","컨테이너 이미지로 Runtime 생성\n자동으로 V1 버전 및 DEFAULT 엔드포인트 생성"),
    (C_PURPLE,"4","에이전트 호출","세션 ID 생성 후 InvokeAgentRuntime 호출\nWebSocket 스트리밍 지원"),
]
for i,(clr,num,t2,d2) in enumerate(steps):
    ty=Inches(2.35)+i*Inches(1.1)
    rc(s,Inches(6.9),ty,Inches(0.45),Inches(0.45),clr)
    tx(s,num,Inches(6.9),ty+Inches(0.05),Inches(0.45),Inches(0.38),
       sz=Pt(16),bold=True,color=C_BG,align=PP_ALIGN.CENTER)
    tx(s,t2,Inches(7.5),ty,Inches(5.2),Inches(0.38),
       sz=Pt(13),bold=True,color=clr)
    tx(s,d2,Inches(7.5),ty+Inches(0.4),Inches(5.2),Inches(0.6),
       sz=Pt(11),color=C_LIGHT)
snum(s,5)

# ══════════════════════════════════════════════════════════════════
# S6 Gateway & Identity
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"04  GATEWAY & IDENTITY")
stitle(s,"Gateway & Identity","도구 연결과 보안 인증·인가의 완전 자동화")
rc(s,Inches(0.45),Inches(1.75),Inches(6.0),Inches(5.3),C_CARD,lc=C_ACCENT2,lw=Pt(1.2))
rc(s,Inches(0.45),Inches(1.75),Inches(6.0),Inches(0.5),C_ACCENT2)
tx(s,"  🔗  AgentCore Gateway",Inches(0.45),Inches(1.75),Inches(6.0),Inches(0.5),
   sz=Pt(15),bold=True,color=C_BG)
gw_items=[
    ("REST API → MCP 변환","OpenAPI 스펙·Smithy 모델 지원\n기존 API를 MCP 서버로 즉시 변환"),
    ("Lambda 함수 연결","서버리스 함수를 에이전트 도구로 등록\n스키마 정의만으로 즉시 사용"),
    ("시맨틱 도구 검색","자연어로 필요한 도구 자동 탐색\n컨텍스트 크기 최적화"),
    ("외부 MCP 서버 연결","기존 MCP 서버 직접 연결 지원\nSlack·Salesforce·GitHub 연동"),
]
for i,(t2,d2) in enumerate(gw_items):
    ty=Inches(2.4)+i*Inches(1.1)
    rc(s,Inches(0.6),ty+Inches(0.1),Inches(0.08),Inches(0.7),C_ACCENT2)
    tx(s,t2,Inches(0.85),ty+Inches(0.05),Inches(5.4),Inches(0.38),
       sz=Pt(13),bold=True,color=C_ACCENT2)
    tx(s,d2,Inches(0.85),ty+Inches(0.45),Inches(5.4),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(5.3),C_CARD,lc=C_GREEN,lw=Pt(1.2))
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.5),C_GREEN)
tx(s,"  🔐  AgentCore Identity",Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.5),
   sz=Pt(15),bold=True,color=C_BG)
id_items=[
    ("OAuth 2.0 / OIDC 표준","업계 표준 인증 프로토콜 완전 지원\n기존 IdP와 원활한 통합"),
    ("2LO · 3LO 지원","클라이언트 자격증명 그랜트 (2LO)\n인가 코드 그랜트 (3LO) 지원"),
    ("최소 권한 접근 제어","동적 사용자 컨텍스트 기반 권한 결정\n에이전트별 스코프 접근 제어"),
    ("커스텀 클레임 지원","멀티테넌트 환경 강화 인증 규칙\n사용자 정의 인증 로직 적용"),
]
for i,(t2,d2) in enumerate(id_items):
    ty=Inches(2.4)+i*Inches(1.1)
    rc(s,Inches(7.0),ty+Inches(0.1),Inches(0.08),Inches(0.7),C_GREEN)
    tx(s,t2,Inches(7.25),ty+Inches(0.05),Inches(5.4),Inches(0.38),
       sz=Pt(13),bold=True,color=C_GREEN)
    tx(s,d2,Inches(7.25),ty+Inches(0.45),Inches(5.4),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
snum(s,6)

# ══════════════════════════════════════════════════════════════════
# S7 Memory & Observability
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"05  MEMORY & OBSERVABILITY")
stitle(s,"Memory & Observability","지능적 기억과 완전한 가시성")
rc(s,Inches(0.45),Inches(1.75),Inches(6.0),Inches(5.3),C_CARD,lc=C_PURPLE,lw=Pt(1.2))
rc(s,Inches(0.45),Inches(1.75),Inches(6.0),Inches(0.5),C_PURPLE)
tx(s,"  🧠  AgentCore Memory",Inches(0.45),Inches(1.75),Inches(6.0),Inches(0.5),
   sz=Pt(15),bold=True,color=C_WHITE)
mem_items=[
    ("단기 메모리","현재 대화 세션 내 컨텍스트 유지\n실시간 상호작용 최적화"),
    ("장기 메모리","사용자별 과거 상호작용 영구 저장\n시간이 지날수록 개인화 향상"),
    ("에피소딕 메모리 (신규)","경험 기반 학습으로 지식 축적\n인간과 유사한 상호작용 구현"),
    ("자체 관리 전략","메모리 추출·통합 파이프라인 완전 제어\n커스텀 메모리 로직 구현 가능"),
]
for i,(t2,d2) in enumerate(mem_items):
    ty=Inches(2.4)+i*Inches(1.1)
    rc(s,Inches(0.6),ty+Inches(0.1),Inches(0.08),Inches(0.7),C_PURPLE)
    tx(s,t2,Inches(0.85),ty+Inches(0.05),Inches(5.4),Inches(0.38),
       sz=Pt(13),bold=True,color=C_PURPLE)
    tx(s,d2,Inches(0.85),ty+Inches(0.45),Inches(5.4),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(5.3),C_CARD,lc=C_TEAL,lw=Pt(1.2))
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.5),C_TEAL)
tx(s,"  📊  AgentCore Observability",Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.5),
   sz=Pt(15),bold=True,color=C_BG)
obs_items=[
    ("OTel 호환 모니터링","OpenTelemetry 표준 완전 지원\n메트릭·스팬·로그 통합 수집"),
    ("CloudWatch 대시보드","에이전트·게이트웨이·메모리 통합 대시보드\n임계값 기반 알림 설정"),
    ("외부 도구 연동","Dynatrace·Datadog·Arize Phoenix\nLangSmith·Langfuse 연동 지원"),
    ("단계별 시각화","에이전트 결정 과정 단계별 추적\n오류 발생 지점 즉시 파악"),
]
for i,(t2,d2) in enumerate(obs_items):
    ty=Inches(2.4)+i*Inches(1.1)
    rc(s,Inches(7.0),ty+Inches(0.1),Inches(0.08),Inches(0.7),C_TEAL)
    tx(s,t2,Inches(7.25),ty+Inches(0.05),Inches(5.4),Inches(0.38),
       sz=Pt(13),bold=True,color=C_TEAL)
    tx(s,d2,Inches(7.25),ty+Inches(0.45),Inches(5.4),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
snum(s,7)

# ══════════════════════════════════════════════════════════════════
# S8 최신 기능 (2025.12)
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"06  LATEST FEATURES")
stitle(s,"최신 기능 (2025년 12월)","Policy · Evaluations · 에피소딕 메모리 · 양방향 스트리밍")
new_features=[
    (C_RED,"🛡️","AgentCore Policy (Preview)",
     "Cedar 기반 정책 언어로 에이전트 행동 규칙 정의\n역할 기반(RBAC) · 값 기반 · 필드 존재 여부 제어\n멀티테넌트 환경에서 세밀한 권한 관리 가능"),
    (C_GREEN,"📈","AgentCore Evaluations (Preview)",
     "에이전트 품질 실시간 모니터링 완전 관리형 서비스\n정확성·안전성·목표 달성률·커스텀 지표 측정\nCloudWatch 통합 대시보드 · 임계값 알림"),
    (C_PURPLE,"🧠","에피소딕 메모리 (Memory)",
     "경험 기반 학습으로 시간이 지날수록 지식 축적\n인간과 유사한 자연스러운 상호작용 구현\n과거 분석 결과를 활용한 더 지능적인 인사이트"),
    (C_ACCENT2,"📡","양방향 스트리밍 (Runtime)",
     "동시 듣기·응답 처리로 자연스러운 대화 구현\n중간 인터럽트 및 컨텍스트 변경 실시간 처리\n음성 에이전트 등 고급 대화형 AI 구현 지원"),
]
for i,(clr,icon,t2,d2) in enumerate(new_features):
    col=i%2; row=i//2
    lx=Inches(0.45)+col*Inches(6.4)
    ty=Inches(1.75)+row*Inches(2.55)
    rc(s,lx,ty,Inches(6.1),Inches(2.35),C_CARD,lc=clr,lw=Pt(1.5))
    rc(s,lx,ty,Inches(6.1),Inches(0.5),clr)
    tx(s,f"{icon}  {t2}",lx+Inches(0.15),ty+Inches(0.07),Inches(5.8),Inches(0.4),
       sz=Pt(14),bold=True,color=C_BG)
    tx(s,d2,lx+Inches(0.2),ty+Inches(0.65),Inches(5.7),Inches(1.55),
       sz=Pt(12),color=C_LIGHT)
snum(s,8)

# ══════════════════════════════════════════════════════════════════
# S9 지원 프레임워크 & 모델 + 엔터프라이즈 보안
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"07 · 08  FRAMEWORKS & SECURITY")
stitle(s,"지원 프레임워크 · 모델 & 엔터프라이즈 보안","오픈 아키텍처 전략과 엔터프라이즈급 보안 기능")

# 프레임워크 섹션
rc(s,Inches(0.45),Inches(1.75),Inches(6.0),Inches(2.5),C_CARD,lc=C_ACCENT,lw=Pt(1.2))
rc(s,Inches(0.45),Inches(1.75),Inches(6.0),Inches(0.45),C_ACCENT)
tx(s,"  🌐  지원 프레임워크",Inches(0.45),Inches(1.75),Inches(6.0),Inches(0.45),
   sz=Pt(13),bold=True,color=C_BG)
frameworks=["LangGraph","CrewAI","LlamaIndex","Google ADK","OpenAI Agents SDK","Strands Agents"]
for i,fw in enumerate(frameworks):
    col=i%3; row=i//3
    lx=Inches(0.65)+col*Inches(1.9)
    ty=Inches(2.35)+row*Inches(0.65)
    rc(s,lx,ty,Inches(1.75),Inches(0.5),C_CARD2,lc=C_ACCENT,lw=Pt(0.8))
    tx(s,fw,lx,ty+Inches(0.08),Inches(1.75),Inches(0.38),
       sz=Pt(11),bold=True,color=C_ACCENT,align=PP_ALIGN.CENTER)

# 모델 섹션
rc(s,Inches(0.45),Inches(4.45),Inches(6.0),Inches(2.1),C_CARD,lc=C_ACCENT2,lw=Pt(1.2))
rc(s,Inches(0.45),Inches(4.45),Inches(6.0),Inches(0.45),C_ACCENT2)
tx(s,"  🤖  지원 모델",Inches(0.45),Inches(4.45),Inches(6.0),Inches(0.45),
   sz=Pt(13),bold=True,color=C_BG)
tx(s,"Amazon Bedrock 내 모든 파운데이션 모델 + 외부 모델 완전 지원\n특정 모델에 종속되지 않는 오픈 아키텍처 채택",
   Inches(0.65),Inches(5.0),Inches(5.6),Inches(1.3),sz=Pt(12),color=C_LIGHT)

# 보안 섹션
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(4.8),C_CARD,lc=C_GREEN,lw=Pt(1.2))
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.45),C_GREEN)
tx(s,"  🔒  엔터프라이즈 보안",Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.45),
   sz=Pt(13),bold=True,color=C_BG)
sec_items=[
    (C_GREEN,"VPC 지원","Virtual Private Cloud 완전 지원\n격리된 네트워크 환경에서 에이전트 실행"),
    (C_ACCENT2,"AWS PrivateLink","인터넷 노출 없이 AWS 서비스 연결\n프라이빗 엔드포인트 통한 안전한 통신"),
    (C_ACCENT,"CloudFormation","인프라 코드(IaC)로 자동화 배포\n반복 가능한 에이전트 인프라 구성"),
    (C_PURPLE,"리소스 태깅","비용 추적 및 거버넌스 관리\n팀·프로젝트별 리소스 분류"),
]
for i,(clr,t2,d2) in enumerate(sec_items):
    ty=Inches(2.35)+i*Inches(1.05)
    rc(s,Inches(7.0),ty,Inches(0.4),Inches(0.4),clr)
    tx(s,t2,Inches(7.55),ty,Inches(5.1),Inches(0.38),
       sz=Pt(13),bold=True,color=clr)
    tx(s,d2,Inches(7.55),ty+Inches(0.42),Inches(5.1),Inches(0.55),
       sz=Pt(11),color=C_LIGHT)
snum(s,9)

# ══════════════════════════════════════════════════════════════════
# S10 활용 사례 & 요약 (마지막 슬라이드)
# ══════════════════════════════════════════════════════════════════
s=ns(); bg(s); topbar(s)
slabel(s,"09  USE CASES & SUMMARY")
stitle(s,"활용 사례 & 핵심 요약","AgentCore로 실현하는 엔터프라이즈 AI 에이전트")

use_cases=[
    (C_ACCENT,"🏢","기업 내부 AI 플랫폼","팀별 에이전트 배포\n공유 메모리 스토어\n거버넌스 중앙화"),
    (C_ACCENT2,"🤝","멀티 에이전트 오케스트레이션","A2A 프로토콜 기반\n에이전트 간 협력\n복잡한 워크플로 자동화"),
    (C_GREEN,"🏥","규제 산업 (의료·금융)","엄격한 보안 기준 유지\nPolicy 기반 행동 제어\n감사 추적 완전 지원"),
    (C_PURPLE,"🎙️","음성 에이전트","양방향 스트리밍 활용\n실시간 대화형 AI\n자연스러운 인터럽트 처리"),
]
for i,(clr,icon,t2,d2) in enumerate(use_cases):
    col=i%2; row=i//2
    lx=Inches(0.45)+col*Inches(3.15)
    ty=Inches(1.75)+row*Inches(2.5)
    rc(s,lx,ty,Inches(2.9),Inches(2.3),C_CARD,lc=clr,lw=Pt(1.2))
    tx(s,icon,lx+Inches(0.15),ty+Inches(0.15),Inches(0.55),Inches(0.55),sz=Pt(22))
    tx(s,t2,lx+Inches(0.15),ty+Inches(0.75),Inches(2.6),Inches(0.45),
       sz=Pt(13),bold=True,color=clr)
    tx(s,d2,lx+Inches(0.15),ty+Inches(1.2),Inches(2.6),Inches(0.95),
       sz=Pt(11),color=C_LIGHT)

# 오른쪽 핵심 요약
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(5.3),C_CARD,lc=C_ACCENT,lw=Pt(1.5))
rc(s,Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.5),C_ACCENT)
tx(s,"  ✅  핵심 가치 요약",Inches(6.85),Inches(1.75),Inches(6.0),Inches(0.5),
   sz=Pt(14),bold=True,color=C_BG)
summary=[
    "🚀  어떤 프레임워크·모델도 자유롭게 사용",
    "🔒  엔터프라이즈급 보안 (VPC·IAM·OAuth)",
    "⚡  최대 8시간 장시간 실행 지원",
    "🧠  에피소딕 메모리로 학습하는 에이전트",
    "🔗  기존 API·Lambda를 즉시 MCP 도구로",
    "📊  CloudWatch 기반 통합 모니터링",
    "🛡️  Policy로 에이전트 행동 규칙 정의",
    "📈  Evaluations로 품질 실시간 측정",
    "🌐  인프라 관리 Zero – 비즈니스 로직에 집중",
]
for i,line in enumerate(summary):
    ty=Inches(2.4)+i*Inches(0.52)
    tx(s,line,Inches(7.05),ty,Inches(5.6),Inches(0.48),
       sz=Pt(12),color=C_LIGHT if i%2==0 else C_WHITE)

# 하단 태그라인
rc(s,Inches(6.85),H-Inches(0.65),Inches(6.0),Inches(0.45),C_ACCENT)
tx(s,"Build · Deploy · Operate  –  Any Agent, Any Scale",
   Inches(6.85),H-Inches(0.65),Inches(6.0),Inches(0.45),
   sz=Pt(12),bold=True,color=C_BG,align=PP_ALIGN.CENTER)

snum(s,10,total=10)

# ══════════════════════════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════════════════════════
out_path = "artifacts/AgentCore_Overview.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}")
